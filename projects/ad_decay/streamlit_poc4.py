#!/usr/bin/env python
# coding: utf-8

# In[ ]:


import streamlit as st
import pandas as pd
import plotly.express as px
import numpy as np
from openai import OpenAI
import re

# -----------------------------
# CONFIG
# -----------------------------
st.set_page_config(layout="wide")
st.title("🎯 Creative Analytics Dashboard")

# -----------------------------
# OPENAI SETUP
# -----------------------------
client = None
if "OPENAI_API_KEY" in st.secrets:
    client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

# -----------------------------
# FILE UPLOAD
# -----------------------------
uploaded_file = st.sidebar.file_uploader("Upload Excel file", type=["xlsx"])

if uploaded_file is None:
    st.warning("Please upload an Excel file.")
    st.stop()

df = pd.read_excel(uploaded_file, engine="openpyxl")
df.columns = df.columns.str.strip()

df["p_date"] = pd.to_datetime(df["p_date"], errors="coerce")
df["Year"] = df["p_date"].dt.year

for col in ["Average Watch Time per Video View", "Cost (USD)", "Purchase", "Checkout"]:
    if col in df.columns:
        df[col] = pd.to_numeric(df[col], errors="coerce")

df["Ad Identifier"] = df["Ad Name"] if "Ad Name" in df.columns else df["Video Material ID"].astype(str)

# -----------------------------
# FILTERS
# -----------------------------
st.sidebar.header("Filters")

year_options = df["Year"].dropna().astype(int).sort_values().unique().tolist()
year_filter = st.sidebar.selectbox(
    "Year",
    ["All"] + year_options
)

country_filter = st.sidebar.selectbox(
    "Ad Country Code",
    ["All"] + sorted(df["Ad Country Code"].dropna().unique())
)
campaign_filter = st.sidebar.selectbox(
    "Campaign Name",
    ["All"] + sorted(df["Campaign Name"].dropna().unique())
)
ad_filter = st.sidebar.selectbox(
    "Ad Name",
    ["All"] + sorted(df["Ad Identifier"].dropna().unique())
)

filtered_df = df.copy()
if year_filter != "All":
    filtered_df = filtered_df[filtered_df["Year"] == year_filter]
if country_filter != "All":
    filtered_df = filtered_df[filtered_df["Ad Country Code"] == country_filter]
if campaign_filter != "All":
    filtered_df = filtered_df[filtered_df["Campaign Name"] == campaign_filter]
if ad_filter != "All":
    filtered_df = filtered_df[filtered_df["Ad Identifier"] == ad_filter]

# -----------------------------
# CONVERSIONS LOGIC
# -----------------------------
def get_conversions(row):
    if row.get("Ad Country Code") == "United States(US)":
        return row.get("Purchase", 0)
    elif row.get("Ad Country Code") == "Australia(AU)":
        return row.get("Checkout", 0)
    return 0

filtered_df["Conversions"] = filtered_df.apply(get_conversions, axis=1)
filtered_df["Conversions"] = filtered_df["Conversions"].fillna(0)

# -----------------------------
# AGGREGATION
# -----------------------------
analysis_level = st.radio("Breakdown", ["Campaign Name", "Ad Name"], horizontal=True)
group_col = "Campaign Name" if analysis_level == "Campaign Name" else "Ad Identifier"

agg_df = filtered_df.groupby(group_col, as_index=False).agg(
    avg_watch_time=("Average Watch Time per Video View", "mean"),
    creative_count=("Ad Identifier", "nunique"),
    cost_usd=("Cost (USD)", "sum"),
    conversions=("Conversions", "sum")
)
agg_df["CPA"] = agg_df.apply(
    lambda row: row["cost_usd"] / row["conversions"] if row["conversions"] > 0 else np.nan,
    axis=1
)

# -----------------------------
# TikTok-style multi-color palette for bubble charts
# -----------------------------
TIKTOK_COLORS = [
    "#FE2C55",
    "#25F4EE",
    "#FF6B35",
    "#A855F7",
    "#F59E0B",
    "#10B981",
    "#F472B6",
    "#818CF8",
    "#34D399",
    "#FB923C",
    "#22D3EE",
    "#FBBF24",
]

QUADRANT_COLORS = {
    "🔴 High Decay / High CPA": "#FE2C55",
    "🟠 High Decay / Low CPA":  "#FF6B35",
    "🟡 Low Decay / High CPA":  "#F59E0B",
    "🟢 Low Decay / Low CPA":   "#10B981",
}

FATIGUE_QUADRANT_COLORS = {
    "🔴 High Fatigue / High CPA": "#FE2C55",
    "🟠 High Fatigue / Low CPA": "#FF6B35",
    "🟡 Low Fatigue / High CPA": "#F59E0B",
    "🟢 Low Fatigue / Low CPA": "#10B981",
}

# -----------------------------
# WATCH TIME BUBBLE
# -----------------------------
st.subheader("📊 Watch Time vs Creative Count")

y_col = "creative_count" if analysis_level == "Campaign Name" else "cost_usd"

legend_order = (
    agg_df.sort_values("avg_watch_time", ascending=False)[group_col]
    .dropna()
    .astype(str)
    .tolist()
)

plot_df = agg_df.copy()
plot_df[group_col] = pd.Categorical(
    plot_df[group_col].astype(str),
    categories=legend_order,
    ordered=True
)

fig = px.scatter(
    plot_df,
    x="avg_watch_time",
    y=y_col,
    size="creative_count",
    color=group_col,
    category_orders={group_col: legend_order},
    color_discrete_sequence=TIKTOK_COLORS,
    hover_data=["CPA", "cost_usd", "conversions"]
)

fig.update_layout(
    legend=dict(traceorder="normal")
)

st.plotly_chart(fig, use_container_width=True)

# -----------------------------
# DECAY SLOPE
# -----------------------------
def compute_slope(group):
    ts = group.groupby("p_date")["Average Watch Time per Video View"].mean()
    if len(ts) > 2:
        return np.polyfit(range(len(ts)), ts.values, 1)[0]
    return np.nan

slope_df = filtered_df.groupby(group_col).apply(compute_slope).reset_index()
slope_df.columns = [group_col, "decay_slope"]
slope_df = slope_df.dropna()

# -----------------------------
# WATCH TIME RANKED DROPDOWN
# -----------------------------
rank_df = agg_df.sort_values("avg_watch_time", ascending=False)
sorted_groups = rank_df[group_col].tolist()

selected_value = st.selectbox(
    "Select group (ranked by highest Watch Time)",
    ["None"] + sorted_groups
)
if selected_value == "None":
    selected_value = None

# -----------------------------
# TIME SERIES WITH DECAY LINE
# -----------------------------
st.subheader("📉 Watch Time Over Time (with Decay Trend)")

ts_df = filtered_df[filtered_df[group_col] == selected_value] if selected_value else filtered_df

ts_agg = ts_df.groupby("p_date", as_index=False).agg(
    avg_watch=("Average Watch Time per Video View", "mean")
).sort_values("p_date").reset_index(drop=True)
ts_agg["t"] = np.arange(len(ts_agg))

if len(ts_agg) > 1:
    slope, intercept = np.polyfit(ts_agg["t"], ts_agg["avg_watch"], 1)
    ts_agg["trend"] = intercept + slope * ts_agg["t"]
else:
    slope = 0
    ts_agg["trend"] = ts_agg["avg_watch"]

fig_ts = px.line(
    ts_agg,
    x="p_date",
    y=["avg_watch", "trend"],
    title="Watch Time with Decay Trend",
    color_discrete_sequence=["#25F4EE", "#FE2C55"],
)
st.plotly_chart(fig_ts, use_container_width=True)
st.metric("Decay Slope", round(slope, 4))

if len(ts_agg) > 2:
    slope, _ = np.polyfit(np.arange(len(ts_agg)), ts_agg["avg_watch"], 1)
else:
    slope = 0

# -----------------------------
# MERGE FOR QUADRANTS
# -----------------------------
quad_df = agg_df.merge(slope_df, on=group_col, how="left")
quad_df = quad_df.dropna(subset=["decay_slope", "CPA"])

slope_median = quad_df["decay_slope"].median()
cpa_median = quad_df["CPA"].median()

def classify(row):
    if row["decay_slope"] <= slope_median and row["CPA"] >= cpa_median:
        return "🔴 High Decay / High CPA"
    elif row["decay_slope"] <= slope_median and row["CPA"] < cpa_median:
        return "🟠 High Decay / Low CPA"
    elif row["decay_slope"] > slope_median and row["CPA"] >= cpa_median:
        return "🟡 Low Decay / High CPA"
    else:
        return "🟢 Low Decay / Low CPA"

quad_df["Quadrant"] = quad_df.apply(classify, axis=1)

# -----------------------------
# TABLE
# -----------------------------
st.subheader("📋 Data (Quadrant Labeled)")
st.dataframe(
    quad_df[[group_col, "avg_watch_time", "creative_count", "cost_usd",
             "conversions", "CPA", "decay_slope", "Quadrant"]],
    use_container_width=True
)

# -----------------------------
# PERCENTILE FATIGUE (QUADRANTS)
# -----------------------------
st.subheader("📈 Fatigue Percentile vs CPA")

slope_df["percentile"] = slope_df["decay_slope"].rank(ascending=True, pct=True)
viz_df = agg_df.merge(slope_df, on=group_col, how="left")
viz_df = viz_df.dropna(subset=["percentile", "CPA", "decay_slope"])

percentile_median = viz_df["percentile"].median()
fatigue_cpa_median = viz_df["CPA"].median()

def classify_percentile_quadrant(row):
    if row["percentile"] <= percentile_median and row["CPA"] >= fatigue_cpa_median:
        return "🔴 High Fatigue / High CPA"
    elif row["percentile"] <= percentile_median and row["CPA"] < fatigue_cpa_median:
        return "🟠 High Fatigue / Low CPA"
    elif row["percentile"] > percentile_median and row["CPA"] >= fatigue_cpa_median:
        return "🟡 Low Fatigue / High CPA"
    else:
        return "🟢 Low Fatigue / Low CPA"

viz_df["Fatigue Quadrant"] = viz_df.apply(classify_percentile_quadrant, axis=1)

fig_fatigue = px.scatter(
    viz_df,
    x="percentile",
    y="CPA",
    size="creative_count",
    color="Fatigue Quadrant",
    color_discrete_map=FATIGUE_QUADRANT_COLORS,
    hover_data=[group_col, "decay_slope", "percentile", "CPA"],
    title="Fatigue Percentile vs CPA"
)

fig_fatigue.add_vline(
    x=percentile_median,
    line_dash="dash",
    line_color="#25F4EE",
    line_width=1.5
)
fig_fatigue.add_hline(
    y=fatigue_cpa_median,
    line_dash="dash",
    line_color="#25F4EE",
    line_width=1.5
)

st.plotly_chart(fig_fatigue, use_container_width=True)

# -----------------------------
# QUADRANT CHART
# -----------------------------
st.subheader("🧭 Quadrant Analysis")

fig_quad = px.scatter(
    quad_df,
    x="decay_slope",
    y="CPA",
    size="creative_count",
    color="Quadrant",
    color_discrete_map=QUADRANT_COLORS,
    hover_data=[group_col]
)
fig_quad.add_vline(x=slope_median, line_dash="dash", line_color="#25F4EE", line_width=1.5)
fig_quad.add_hline(y=cpa_median, line_dash="dash", line_color="#25F4EE", line_width=1.5)
st.plotly_chart(fig_quad, use_container_width=True)

# -----------------------------
# LOW DATA DIAGNOSTIC
# -----------------------------
st.subheader("⚠️ Campaigns with Limited Data")

days_df = filtered_df.groupby(group_col)["p_date"].nunique().reset_index()
days_df.columns = [group_col, "active_days"]

diagnostic_df = agg_df.merge(days_df, on=group_col, how="left")
diagnostic_df = diagnostic_df.merge(slope_df, on=group_col, how="left")
diagnostic_df["<2_days"] = diagnostic_df["active_days"] < 2
diagnostic_df["zero_conversions"] = diagnostic_df["conversions"] == 0
diagnostic_df["issue_flag"] = diagnostic_df.apply(
    lambda row: "⚠️ Both" if row["<2_days"] and row["zero_conversions"]
    else ("⚠️ <2 days" if row["<2_days"]
    else ("⚠️ 0 conversions" if row["zero_conversions"]
    else "OK")),
    axis=1
)

problem_df = diagnostic_df[
    (diagnostic_df["<2_days"]) | (diagnostic_df["zero_conversions"])
]

st.write(f"Total groups: {len(diagnostic_df)}")
st.write(f"Problem groups: {len(problem_df)}")

st.dataframe(
    problem_df[[group_col, "active_days", "conversions", "cost_usd",
                "avg_watch_time", "decay_slope", "issue_flag"]]
    .sort_values(["active_days", "conversions"]),
    use_container_width=True
)

# -----------------------------
# AI INSIGHTS HELPERS
# -----------------------------
def split_insights(text):
    """
    Split AI output into separate insight blocks.
    Supports:
    - numbered lists: 1. ..., 2. ...
    - bullet lists: - ..., * ..., • ...
    - paragraph-separated insights
    """
    if not text or not text.strip():
        return []

    text = text.strip()

    numbered = re.split(r"\n(?=\d+[\.\)])", text)
    if len(numbered) > 1:
        return [x.strip() for x in numbered if x.strip()]

    bulleted = re.split(r"\n(?=[\-\*•]\s)", text)
    if len(bulleted) > 1:
        return [x.strip() for x in bulleted if x.strip()]

    paragraphs = [x.strip() for x in text.split("\n\n") if x.strip()]
    return paragraphs if paragraphs else [text]

# -----------------------------
# AI INSIGHTS
# -----------------------------
st.subheader("🧠 AI Insights")

if st.button("Generate Insights"):
    if client is None:
        st.warning("OpenAI not configured")
    else:
        table_text = quad_df.to_string(index=False)
        prompt = f"""
You are analyzing a marketing dataset.

Each row includes Watch time, CPA, Decay slope, and Quadrant:
🔴 High Decay / High CPA | 🟠 High Decay / Low CPA | 🟡 Low Decay / High CPA | 🟢 Low Decay / Low CPA

Use ONLY the table below.
Focus on:
1. quadrant distribution
2. high risk groups (🔴)
3. efficiency opportunities (🟢)
4. actionable recommendations

Return the answer as a numbered list.
Each numbered item must be a standalone insight suitable for one PowerPoint slide.
Keep each insight concise but specific.

DATA:
{table_text}
"""
        response = client.chat.completions.create(
            model="gpt-4.1-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        ai_text = response.choices[0].message.content
        st.session_state["ai_insights"] = ai_text

        insight_list = split_insights(ai_text)
        for i, insight in enumerate(insight_list, start=1):
            st.markdown(f"**Insight {i}**")
            st.write(insight)

# =============================================================
# EXPORT — PROFESSIONAL PPTX (TikTok For Business dark theme)
# One chart per slide
# =============================================================
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import tempfile, os
import plotly.io as pio
import plotly.graph_objects as go

st.subheader("📤 Export Professional PPT Report")

BG_BLACK = RGBColor(0x01, 0x01, 0x01)
BG_DARK = RGBColor(0x12, 0x12, 0x12)
ACCENT_RED = RGBColor(0xFE, 0x2C, 0x55)
ACCENT_CYAN = RGBColor(0x25, 0xF4, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xAA, 0xAA, 0xAA)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)

PPT_PALETTE = [
    "#FE2C55", "#25F4EE", "#FF6B35", "#A855F7",
    "#F59E0B", "#10B981", "#F472B6", "#818CF8",
    "#34D399", "#FB923C", "#22D3EE", "#FBBF24",
]

def style_fig_for_export(f, multi_color=False):
    f.update_layout(
        template="plotly_dark",
        font=dict(family="Arial Black, Arial, sans-serif", size=13, color="#FFFFFF"),
        paper_bgcolor="#121212",
        plot_bgcolor="#121212",
        margin=dict(l=70, r=40, t=55, b=65),
        legend=dict(
            bgcolor="rgba(18,18,18,0.9)",
            bordercolor="#333333",
            borderwidth=1,
            font=dict(color="#FFFFFF", size=11),
        ),
        title_font=dict(family="Arial Black, Arial", size=16, color="#FFFFFF"),
    )
    f.update_xaxes(
        gridcolor="#2a2a2a", linecolor="#333333",
        tickfont=dict(color="#AAAAAA"), title_font=dict(color="#AAAAAA"),
        zerolinecolor="#333333",
    )
    f.update_yaxes(
        gridcolor="#2a2a2a", linecolor="#333333",
        tickfont=dict(color="#AAAAAA"), title_font=dict(color="#AAAAAA"),
        zerolinecolor="#333333",
    )
    if multi_color:
        for i, trace in enumerate(f.data):
            if hasattr(trace, "marker"):
                trace.marker.color = PPT_PALETTE[i % len(PPT_PALETTE)]
    return f

def save_chart(f, path, w=1600, h=820, multi_color=False):
    import copy
    pio.write_image(
        style_fig_for_export(copy.deepcopy(f), multi_color=multi_color),
        path, format="png", width=w, height=h, scale=2
    )

def save_table(data_df, path):
    row_colors = ["#1a1a1a" if i % 2 == 0 else "#121212" for i in range(len(data_df))]
    tbl = go.Figure(data=[go.Table(
        header=dict(
            values=[f"<b>{c}</b>" for c in data_df.columns],
            fill_color="#FE2C55",
            font=dict(color="white", size=12, family="Arial Black, Arial"),
            align="left", height=34,
            line=dict(color="#333333", width=1),
        ),
        cells=dict(
            values=[data_df[c] for c in data_df.columns],
            fill_color=[row_colors],
            font=dict(color="#FFFFFF", size=11, family="Arial, sans-serif"),
            align="left", height=30,
            line=dict(color="#333333", width=1),
        ),
    )])
    tbl.update_layout(
        margin=dict(l=10, r=10, t=10, b=10),
        paper_bgcolor="#121212",
    )
    tbl.write_image(path, width=1600, height=min(90 + len(data_df) * 34, 900), scale=2)

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_textbox(slide, text, x, y, w, h, size,
                bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                italic=False, font="Arial"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb

def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG_BLACK
    bar.line.fill.background()

    sep = slide.shapes.add_shape(1, Inches(0), Inches(1.06), SLIDE_W, Inches(0.04))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ACCENT_RED
    sep.line.fill.background()

    add_textbox(
        slide, title,
        Inches(0.45), Inches(0.08), Inches(12.4), Inches(0.65),
        size=28, bold=True, color=WHITE,
        align=PP_ALIGN.LEFT, font="Arial Black"
    )

    if subtitle:
        add_textbox(
            slide, subtitle,
            Inches(0.45), Inches(0.72), Inches(12.4), Inches(0.34),
            size=11, color=MUTED, align=PP_ALIGN.LEFT, font="Arial"
        )

def add_chart_slide(prs, title, subtitle, image_path, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_title_bar(slide, title, subtitle)

    img_w = Inches(12.33)
    img_left = (SLIDE_W - img_w) / 2
    slide.shapes.add_picture(image_path, img_left, Inches(1.18), width=img_w)

    if note:
        add_textbox(
            slide, note,
            MARGIN, Inches(6.92), Inches(12.33), Inches(0.45),
            size=10, color=MUTED, italic=True, font="Arial"
        )

def add_title_slide(prs, title, subtitle, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_BLACK)

    words = title.rsplit(" ", 1)
    main_t = words[0] + " " if len(words) > 1 else title
    red_t = words[1] if len(words) > 1 else ""

    tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.33), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    r1 = p.add_run()
    r1.text = main_t
    r1.font.size = Pt(44)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r1.font.name = "Arial Black"

    if red_t:
        r2 = p.add_run()
        r2.text = red_t
        r2.font.size = Pt(44)
        r2.font.bold = True
        r2.font.color.rgb = ACCENT_RED
        r2.font.name = "Arial Black"

    add_textbox(
        slide, subtitle,
        Inches(1.5), Inches(3.75), Inches(10.33), Inches(0.8),
        size=17, color=MUTED, align=PP_ALIGN.CENTER, font="Arial"
    )

    dot1 = slide.shapes.add_shape(1, Inches(6.4), Inches(4.65), Inches(0.14), Inches(0.14))
    dot1.fill.solid()
    dot1.fill.fore_color.rgb = ACCENT_RED
    dot1.line.fill.background()

    dot2 = slide.shapes.add_shape(1, Inches(6.62), Inches(4.65), Inches(0.14), Inches(0.14))
    dot2.fill.solid()
    dot2.fill.fore_color.rgb = ACCENT_CYAN
    dot2.line.fill.background()

    add_textbox(
        slide, meta,
        Inches(1.5), Inches(6.6), Inches(10.33), Inches(0.5),
        size=11, color=RGBColor(0x55, 0x55, 0x55),
        italic=True, align=PP_ALIGN.CENTER, font="Arial"
    )

def add_insight_slide(prs, insight_text, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_title_bar(
        slide,
        f"AI Insight {idx}",
        f"Generated analysis based on quadrant data ({idx}/{total})"
    )
    add_textbox(
        slide,
        insight_text.strip(),
        MARGIN, Inches(1.25), Inches(12.33), Inches(5.8),
        size=18, color=WHITE, font="Arial"
    )

def create_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    tmp = tempfile.mkdtemp()

    p_bubble = os.path.join(tmp, "bubble.png")
    p_ts = os.path.join(tmp, "ts.png")
    p_fatigue = os.path.join(tmp, "fatigue.png")
    p_quad = os.path.join(tmp, "quad.png")
    p_tbl = os.path.join(tmp, "tbl_quad.png")

    save_chart(fig, p_bubble, multi_color=True)
    save_chart(fig_ts, p_ts, multi_color=False)
    save_chart(fig_fatigue, p_fatigue, multi_color=False)
    save_chart(fig_quad, p_quad, multi_color=False)

    save_table(
        quad_df[[group_col, "avg_watch_time", "cost_usd",
                 "conversions", "CPA", "decay_slope", "Quadrant"]].head(20),
        p_tbl
    )

    p_low = None
    try:
        if not problem_df.empty:
            p_low = os.path.join(tmp, "tbl_low.png")
            save_table(
                problem_df[[group_col, "active_days", "conversions",
                            "cost_usd", "avg_watch_time", "issue_flag"]].head(20),
                p_low
            )
    except Exception:
        pass

    meta = f"Year: {year_filter}  |  Country: {country_filter}  |  Campaign: {campaign_filter}  |  Ad: {ad_filter}"

    add_title_slide(
        prs,
        title="Creative Analytics Report",
        subtitle=f"Performance · Fatigue · Efficiency  |  {analysis_level} View",
        meta=meta,
    )

    add_chart_slide(
        prs,
        title="Watch Time vs Creative Count",
        subtitle="Bubble size = number of creatives  |  Legend ordered by highest watch time",
        image_path=p_bubble,
        note="Higher-watch-time groups appear first in the legend. Hover for CPA, cost and conversions.",
    )

    add_chart_slide(
        prs,
        title="Watch Time Over Time",
        subtitle=f"Linear decay trend overlay  |  Slope: {round(slope, 4)}",
        image_path=p_ts,
        note="A negative slope signals audience fatigue. Cyan = actual watch time, Red = linear trend.",
    )

    add_chart_slide(
        prs,
        title="Fatigue Percentile vs CPA",
        subtitle="Quadrant grouping based on percentile and CPA medians",
        image_path=p_fatigue,
        note="Upper-left and upper-right regions highlight higher CPA groups. Dashed lines mark median thresholds.",
    )

    add_chart_slide(
        prs,
        title="Quadrant Analysis",
        subtitle="🔴 High Decay/High CPA   🟠 High Decay/Low CPA   🟡 Low Decay/High CPA   🟢 Low Decay/Low CPA",
        image_path=p_quad,
        note="Dashed lines mark median thresholds. 🟢 = healthy efficiency. 🔴 = immediate action required.",
    )

    add_chart_slide(
        prs,
        title="Quadrant Data Table",
        subtitle=f"Top 20 {analysis_level.lower()}s — full performance breakdown",
        image_path=p_tbl,
    )

    if p_low:
        add_chart_slide(
            prs,
            title="⚠️ Low Data — Action Required",
            subtitle="Groups with fewer than 2 active days or zero conversions",
            image_path=p_low,
            note="Decay slopes for these groups may be unreliable. Verify data pipeline or allow more time to accumulate.",
        )

    insights = st.session_state.get("ai_insights")
    if insights:
        insight_list = split_insights(insights)
        for i, insight in enumerate(insight_list, start=1):
            add_insight_slide(prs, insight, i, len(insight_list))

    out = os.path.join(tmp, "creative_analytics_report.pptx")
    prs.save(out)
    return out

# ── Export button ─────────────────────────────────────────────
if st.button("Generate Professional PPT"):
    with st.spinner("Building slides…"):
        try:
            st.session_state["ppt_file"] = create_ppt()
            st.success("✅ PPT ready — click below to download.")
        except Exception as e:
            st.error(f"Export failed: {e}")
            st.exception(e)

if "ppt_file" in st.session_state:
    with open(st.session_state["ppt_file"], "rb") as f:
        st.download_button(
            label="📥 Download PPT",
            data=f,
            file_name="creative_analytics_report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

